#!/usr/bin/env python3
"""
Genera dorsal y postal personalizados en PDF.
Recibe JSON por stdin con: nombre, bib_number, dorsal_template (base64), postal_template (base64)
Imprime JSON con: dorsal_pdf (base64), postal_pdf (base64)
"""
import sys, io, json, base64, os, subprocess, tempfile
from pptx import Presentation

def reemplazar_en_pptx(pptx_bytes, reemplazos):
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in reemplazos.items():
                        placeholder = '{{' + key + '}}'
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def pptx_a_pdf_base64(pptx_bytes, nombre_temp):
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, f'{nombre_temp}.pptx')
        pdf_path = os.path.join(tmpdir, f'{nombre_temp}.pdf')
        with open(pptx_path, 'wb') as f:
            f.write(pptx_bytes)
        subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pdf',
            '--outdir', tmpdir, pptx_path
        ], capture_output=True, timeout=60)
        with open(pdf_path, 'rb') as f:
            return base64.b64encode(f.read()).decode('utf-8')

def main():
    data = json.loads(sys.stdin.read())
    nombre = data['nombre']
    bib_number = str(data['bib_number'])
    dorsal_bytes = base64.b64decode(data['dorsal_template'])
    postal_bytes = base64.b64decode(data['postal_template'])

    dorsal_modificado = reemplazar_en_pptx(dorsal_bytes, {
        'bib_number': bib_number,
        'nombre_atleta': nombre
    })
    dorsal_pdf = pptx_a_pdf_base64(dorsal_modificado, 'dorsal')

    postal_modificado = reemplazar_en_pptx(postal_bytes, {
        'nombre_atleta': nombre
    })
    postal_pdf = pptx_a_pdf_base64(postal_modificado, 'postal')

    print(json.dumps({ 'dorsal_pdf': dorsal_pdf, 'postal_pdf': postal_pdf }))

if __name__ == '__main__':
    main()