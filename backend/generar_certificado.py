#!/usr/bin/env python3
"""
Genera certificado de finisher personalizado en PDF.
Recibe JSON por stdin con: nombre, nombre_desafio, distancia_km, bib_number, fecha_completado, certificado_template (base64)
Imprime JSON con: certificado_pdf (base64)
"""
import sys, io, json, base64, os, subprocess, tempfile
from pptx import Presentation

def reemplazar_en_pptx(pptx_bytes, reemplazos):
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in reemplazos.items():
                        placeholder = '{{' + key + '}}'
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def pptx_a_pdf_base64(pptx_bytes, nombre_temp):
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, f'{nombre_temp}.pptx')
        pdf_path = os.path.join(tmpdir, f'{nombre_temp}.pdf')
        with open(pptx_path, 'wb') as f:
            f.write(pptx_bytes)
        subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pdf',
            '--outdir', tmpdir, pptx_path
        ], capture_output=True, timeout=60)
        with open(pdf_path, 'rb') as f:
            return base64.b64encode(f.read()).decode('utf-8')

def main():
    data = json.loads(sys.stdin.read())
    nombre = data['nombre']
    nombre_desafio = data['nombre_desafio']
    distancia_km = str(data['distancia_km'])
    bib_number = str(data['bib_number'])
    fecha_completado = data['fecha_completado']
    numero_serie = data['numero_serie']
    certificado_bytes = base64.b64decode(data['certificado_template'])

    certificado_modificado = reemplazar_en_pptx(certificado_bytes, {
        'NOMBRE_COMPLETO': nombre,
        'NOMBRE_DESAFIO': nombre_desafio,
        'DISTANCIA_KM': distancia_km,
        'BIB_NUMBER': bib_number,
        'FECHA_COMPLETADO': fecha_completado,
        'NUMERO_SERIE': numero_serie,
    })
    certificado_pdf = pptx_a_pdf_base64(certificado_modificado, 'certificado')

    print(json.dumps({ 'certificado_pdf': certificado_pdf }))

if __name__ == '__main__':
    main()
    